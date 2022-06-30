import json
import os
import ssl
import logging
import re
import shutil

import bson
from abc import ABC, abstractmethod

from docx import Document
from docx.opc.exceptions import PackageNotFoundError
from pptx import Presentation
import fitz

from pymongo.mongo_client import MongoClient
from analysis import CharDictAnalysis, CharDictionary

from gensim.models import Word2Vec
from gensim import matutils  # utility fnc for pickling, common scipy operations etc
import pymorphy2
import numpy as np

import win32com.client as win32
from win32com.client import constants

db_name = 'characteristics'
connection_string = ''
foot_prints_dir = 'footprints'
models_dir = 'models'


class Basic(ABC):
    logger = None

    @abstractmethod
    def __init__(self):
        self.get_logger()

    @abstractmethod
    def get_logger(self):
        logger = logging.getLogger("core")

        if not len(logger.handlers):
            formatter = logging.Formatter(
                '%(asctime)s - %(module)s - %(levelname)s - %(funcName)s: %(lineno)d - %(message)s',
                datefmt="'%H:%M:%S',")

            consoleHandler = logging.StreamHandler()
            consoleHandler.setFormatter(formatter)

            logger.addHandler(consoleHandler)
            logger.setLevel(logging.INFO)

        self.logger = logger

    @abstractmethod
    def connect_db(self, name):
        client = MongoClient(
            connection_string, ssl_cert_reqs=ssl.CERT_NONE)
        return client[name]


class Standard(Basic):
    direction = None

    def __init__(self, direction):
        super().__init__()
        self.direction = direction
        self.db = self.connect_db(db_name)

    # Сформировать модели проф. стандартов
    def form_ps_models(self, ps_code_dir):
        regx = bson.regex.Regex(f'^{ps_code_dir}....')
        dir_ps_codes = self.db['prof_standarts'].find({'code': regx})

        for ps in dir_ps_codes:
            ps_sentences = self.get_sentences_from_code(ps['code'])
            self.logger.debug(f"{ps['code']} | {ps_sentences}")
            model = Word2VecModel(self.direction, ps['code'])
            model.create_and_train_model(ps_sentences)

    # Сформировать модели ФГОС
    def form_fgos_models(self, fgos_code_dir):
        regx = bson.regex.Regex(f'^{fgos_code_dir}......')
        dir_fgos_codes = self.db['fgos'].find({'code': regx})

        for fgos in dir_fgos_codes:
            fgos_sentences = self.get_sentences_from_code(fgos['code'])
            self.logger.debug(f"{fgos['code']} | {fgos_sentences}")
            model = Word2VecModel(self.direction, fgos['code'])
            model.create_and_train_model(fgos_sentences)

    # Сформировать все модели по направлению
    def form_direction_models(self):
        ps_code_dir = None
        fgos_code_dir = None

        if self.direction == 'it':
            ps_code_dir = '06'
            fgos_code_dir = '09'
        elif self.direction == 'ped':
            ps_code_dir = '01'
            fgos_code_dir = '44'

        self.form_ps_models(ps_code_dir)
        self.form_fgos_models(fgos_code_dir)

    # Получить состав предложений по коду ПС/ФГОС
    def get_sentences_from_code(self, code):
        sentences = []

        check_ps = re.match(r'^\d{2}.\d{3}$', code)
        check_fgos = re.match(r'^\d{2}.\d{2}.\d{2}$', code)

        # ПС
        if check_ps is not None:
            tfs = self.db['tfs'].find({'ps_code': code})
            if tfs is None:
                raise Exception(f"Not found ps: {code}")

            for tf in tfs:
                all_s = []
                if 'required_skills' in tf:
                    if tf['required_skills'] is not None:
                        all_s += tf['required_skills']
                if 'required_knowledge' in tf:
                    if tf['required_knowledge'] is not None:
                        all_s += tf['required_knowledge']
                if 'labor_actions' in tf:
                    if tf['labor_actions'] is not None:
                        all_s += tf['labor_actions']

                for s in all_s:
                    s = CharDictionary.text_treatment(s)
                    s = CharDictionary.lemmatize_sentence(s)
                    if len(s) > 0:
                        sentences.append(s)

        # ФГОС
        elif check_fgos is not None:
            fgos = self.db['fgos'].find_one({'code': code})
            if fgos is None:
                raise Exception(f"Not found fgos: {code}")

            # Унив. комп.
            uk_codes = fgos['uk_codes']
            for uk_code in uk_codes:
                uk = self.db['uk'].find_one({'code': uk_code})
                s = CharDictionary.text_treatment(uk['name'])
                s = CharDictionary.lemmatize_sentence(s)
                if len(s) > 0:
                    sentences.append(s)

            # ОПК
            if 'opk_codes' in fgos and fgos['opk_codes'] is not None:
                for opk_code in fgos['opk_codes']:
                    opk = self.db['opk'].find_one({'direction': code[:2], 'code': opk_code})
                    s = CharDictionary.text_treatment(opk['name'])
                    s = CharDictionary.lemmatize_sentence(s)
                    if len(s) > 0:
                        sentences.append(s)

            ps_codes = fgos['ps_codes']
            for ps_code in ps_codes:
                # TODO 01.004 костыль
                if ps_code == '01.004':
                    continue
                tfs = self.db['tfs'].find({'ps_code': ps_code})
                if tfs is None:
                    raise Exception(f"Not found ps: {ps_code}")

                # Повтор кода 123
                for tf in tfs:
                    all_s = []
                    if 'required_skills' in tf:
                        if tf['required_skills'] is not None:
                            all_s += tf['required_skills']
                    if 'required_knowledge' in tf:
                        if tf['required_knowledge'] is not None:
                            all_s += tf['required_knowledge']
                    if 'labor_actions' in tf:
                        if tf['labor_actions'] is not None:
                            all_s += tf['labor_actions']

                    for s in all_s:
                        s = CharDictionary.text_treatment(s)
                        s = CharDictionary.lemmatize_sentence(s)
                        if len(s) > 0:
                            sentences.append(s)

        else:
            raise Exception("Unknown code format")

        return sentences

    def get_st(self):
        if self.direction == "it":
            return self.get_it_st()
        elif self.direction == "ped":
            return self.get_ped_st()

    # Получение проф. стандартов IT направления
    def get_it_st(self):
        # Получем только ПС it направления
        it_stop_ps = ['06.023', '06.040', '06.039', '06.005', '06.043',
                      '06.045', '06.037', '06.038', '06.036', '06.007',
                      '06.030', '06.021', '06.009', '06.002', '06.008',
                      '06.029', '06.027', '06.026', '06.024', '06.020',
                      '06.018', '06.010', '06.006']

        regx = bson.regex.Regex(f'^06....')
        all_ps = self.db['prof_standarts'].find({'code': regx})

        it_st = []
        for ps in all_ps:
            if ps['code'] not in it_stop_ps:
                it_st.append(ps['code'])

        regx = bson.regex.Regex(f'^09......')
        dir_fgos_codes = self.db['fgos'].find({'code': regx})

        for fgos in dir_fgos_codes:
            it_st.append(fgos['code'])

        return it_st

    def get_ped_st(self):
        regx = bson.regex.Regex(f'^01....')
        all_ps = self.db['prof_standarts'].find({'code': regx})

        ped_st = []

        # TODO Костыль на 01.006 (нет в словаре характеристик)
        for ps in all_ps:
            ped_st.append(ps['code'])

        regx = bson.regex.Regex(f'^44......')
        dir_fgos_codes = self.db['fgos'].find({'code': regx})

        for fgos in dir_fgos_codes:
            ped_st.append(fgos['code'])

        return ped_st

    def get_char_dict_words_by_code(self, code):
        dict = self.db['dictionaries'].find_one({'name': self.direction})['content']
        words = None

        for ps_fgos in dict:
            if 'ps_code' in ps_fgos and ps_fgos['ps_code'] == code:
                words = ps_fgos['words']
            elif 'fgos_code' in ps_fgos and ps_fgos['fgos_code'] == code:
                words = ps_fgos['words']

        return words

    def comparison_all_models(self, mode, user=None, foot_prints_count=None, save_as="json"):
        st_comp = self.get_st()

        ps_result = None
        if mode == "all":
            ps_result = {}
        elif mode == "user":
            ps_result = []

        cda = CharDictAnalysis(self.direction)
        for code in st_comp:
            # TODO костыль 006, 007
            if code == '01.006' or code == '01.007':
                continue

            words = self.get_char_dict_words_by_code(code)

            if words is None:
                raise Exception(f"Not found code: {code}")

            if mode == "all":
                ps_result[code] = {
                    'st_name': cda.get_st_name_by_code(code),
                    'comparison': []
                }

                for ps2 in st_comp:
                    code_two = ps2

                    model = Word2VecModel(self.direction, code)
                    model2 = Word2VecModel(self.direction, code_two)
                    model.load_model()
                    model2.load_model()

                    similarities = Word2VecModel.comparison_models(model, model2, words)

                    ps_result[code]['comparison'].append({
                        'st_code': ps2,
                        'st_name': cda.get_st_name_by_code(ps2),
                        'similarities': similarities,
                    })

            elif mode == "user":
                if user is None:
                    raise Exception("User not founded")
                model = Word2VecModel(self.direction, code)
                model2 = Word2VecModel("users", user)
                model.load_model()
                model2.load_model()

                similarities, hits = Word2VecModel.comparison_models(model, model2, words)

                # coincidences = self.find_coincidence(code, hits)

                ps_result.append({
                    'st_code': code,
                    'st_name': CharDictAnalysis(self.direction).get_st_name_by_code(code),
                    'similarities': similarities,
                    'hits': hits,
                    # 'coincidences': coincidences
                })

        if mode == "all":
            for w in ps_result:
                ps_result[w]['comparison'] = sorted(ps_result[w]['comparison'], key=lambda d: d['similarities'])

            with open(f'{self.direction}_comparison.json', 'w+', encoding='utf8') as outfile:
                json.dump(ps_result, outfile, indent=4, ensure_ascii=False)

        elif mode == "user":
            ps_result = sorted(ps_result, key=lambda d: d['similarities'])

            if save_as == "db":
                self.db[f'assessment'].insert_one({
                    'student_login': user,
                    'direction': self.direction,
                    'foot_prints_count': foot_prints_count,
                    'result': ps_result
                })

                self.logger.info("assessment result successfully saved in db")
            elif save_as == "json":
                with open(f'assessment/{self.direction}/{user}_comparison.json', 'w+', encoding='utf8') as outfile:
                    json.dump(ps_result, outfile, indent=4, ensure_ascii=False)

    @staticmethod
    def training_and_comparison_all_users_by_dir(direction, min_count_footprints, save_as="json"):
        path = f"{foot_prints_dir}/{direction}"
        content_dir = os.listdir(path)
        users_logins = []

        for c in content_dir:
            if os.path.isdir(os.path.join(path, c)):
                users_logins.append(c)

        for login in users_logins:
            client = MongoClient(
                connection_string, ssl_cert_reqs=ssl.CERT_NONE)
            db = client[db_name]
            if db[f'assessment'].find_one({'student_login': login}) is not None:
                print(f"User {login} already assessment, skip")
                continue

            footprints = FootPrint(login, direction).get_all_foot_prints_from_dir(min_count_footprints)
            if footprints is None:
                print("Skip little footprints student")
                continue
            Word2VecModel('users', login).create_and_training_user_model(footprints)
            Standard(direction).comparison_all_models("user", login, len(footprints), save_as)

    # TODO изменить
    # def find_coincidence(self, ps_code, hits):
    #     coincidences = []
    #     tfs = self.db['tfs'].find({'ps_code': ps_code})
    #     if tfs is None:
    #         raise Exception(f"Not found ps: {ps_code}")
    #
    #     for tf in tfs:
    #         if 'required_skills' in tf:
    #             if tf['required_skills'] is not None:
    #                 for skill in tf['required_skills']:
    #                     s = CharDictionary.text_treatment(skill)
    #                     s = CharDictionary.lemmatize_sentence(s)
    #                     if len(s) > 0:
    #                         for hit in hits.keys():
    #                             if hit in s and hits[hit]['p'] > 0.0001:
    #                                 coincidences.append({
    #                                     'key': hit,
    #                                     'st': ps_code,
    #                                     'tf': tf['code'],
    #                                     'type': 'skill',
    #                                     'text': skill
    #                                 })
    #
    #         if 'required_knowledge' in tf:
    #             if tf['required_knowledge'] is not None:
    #                 for knowledge in tf['required_knowledge']:
    #                     s = CharDictionary.text_treatment(knowledge)
    #                     s = CharDictionary.lemmatize_sentence(s)
    #                     if len(s) > 0:
    #                         for hit in hits.keys():
    #                             if hit in s and hits[hit]['p'] > 0.0001:
    #                                 coincidences.append({
    #                                     'source': f"PS: {ps_code} | TF: {tf['code']} | Type: Knowledge | Text: {knowledge}"
    #                                 })
    #
    #         if 'labor_actions' in tf:
    #             if tf['labor_actions'] is not None:
    #                 for actions in tf['labor_actions']:
    #                     s = CharDictionary.text_treatment(actions)
    #                     s = CharDictionary.lemmatize_sentence(s)
    #                     if len(s) > 0:
    #                         for hit in hits.keys():
    #                             if hit in s and hits[hit]['p'] > 0.0001:
    #                                 coincidences.append({
    #                                     'source': f"PS: {ps_code} | TF: {tf['code']} | Type: Actions | Text: {actions}"
    #                                 })
    #     return coincidences

    def get_logger(self):
        super().get_logger()

    def connect_db(self, name):
        return super().connect_db(name)


class FootPrint(Basic):
    user = None
    sentences = None

    def __init__(self, user, direction, file=None):
        super().__init__()
        self.user = user
        self.direction = direction
        self.file = file

    # Получить состав предложений цифрового следа
    def get_sentences(self, footprint_type, file):
        sentences = []
        file_path = f"{foot_prints_dir}/{self.direction}/{self.user}/{file}"
        if footprint_type == 'word':
            doc = Document(file_path)
            para = doc.paragraphs
            for p in para:
                if len(p.text) > 0:
                    text = CharDictionary.text_treatment(p.text)
                    sentence = CharDictionary.lemmatize_sentence(text)
                    if len(sentence) > 3:
                        sentences.append(sentence)
        elif footprint_type == 'pp':
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        for s in shape.text.split('\n'):
                            if len(s) > 0:
                                text = CharDictionary.text_treatment(s)
                                sentence = CharDictionary.lemmatize_sentence(text)
                                if len(sentence) > 3:
                                    sentences.append(sentence)
        elif footprint_type == 'pdf':
            with fitz.open(file_path) as doc:
                for page in doc:
                    for s in page.get_text().split('\n'):
                        if len(s) > 0:
                            text = CharDictionary.text_treatment(s)
                            sentence = CharDictionary.lemmatize_sentence(text)
                            if len(sentence) > 3:
                                sentences.append(sentence)
        else:
            raise Exception("Unknown footprint type")
        self.sentences = sentences
        self.logger.info("Sentences successfully generated")

    def get_all_foot_prints_from_dir(self, min_foot_prints_count):
        path = os.path.join(foot_prints_dir, self.direction, self.user)
        # Формируем список из файлов директории path с проверкой на файл (формируем абсолютный путь для этого)
        files = [f for f in os.listdir(path) if os.path.isfile(os.path.join(path, f))]

        if len(files) < min_foot_prints_count:
            return None

        word_formats = ["docx", "doc"]
        pp_formats = ["pptx", "ppt"]

        footprints = []
        for f in files:
            path_file = os.path.join(path, f)
            extension = f.split(".")[-1]
            if extension in word_formats:
                if extension == "doc":
                    FootPrint.save_doc_as_docx(os.path.join(os.getcwd(), path_file))
                    f += "x"
                foot_type = 'word'
            elif extension in pp_formats:
                if extension == "ppt":
                    FootPrint.save_ppt_as_pptx(os.path.join(os.getcwd(), path_file))
                    f += "x"
                foot_type = 'pp'
            elif extension == 'pdf':
                foot_type = 'pdf'
            else:
                os.remove(path_file)
                self.logger.warning("Remove footprint file with undefined format: ", path_file)
                continue

            foot = FootPrint(self.user, self.direction, f)
            try:
                foot.get_sentences(foot_type, f)
            except PackageNotFoundError:
                continue
            footprints.append(foot)

            self.logger.info(f"[{self.user}] File print load and sentences gutted: {f}, type - {foot_type}")

        return footprints

    def get_logger(self):
        super().get_logger()

    def connect_db(self, name):
        return super().connect_db(name)

    @staticmethod
    def save_doc_as_docx(path):
        # Opening MS Word
        word = win32.gencache.EnsureDispatch('Word.Application')
        doc = word.Documents.Open(path)
        doc.Activate()

        # Rename path with .docx
        new_file_abs = os.path.abspath(path)
        new_file_abs = re.sub(r'\.\w+$', '.docx', new_file_abs)

        # Save and Close
        word.ActiveDocument.SaveAs(
            new_file_abs, FileFormat=constants.wdFormatXMLDocument
        )
        doc.Close(False)
        os.remove(path)

    @staticmethod
    def save_ppt_as_pptx(path):
        PptApp = win32.Dispatch("Powerpoint.Application")
        PPtPresentation = PptApp.Presentations.Open(path)
        PPtPresentation.SaveAs(path + "x", 24)
        PptApp.Quit()
        os.remove(path)


# Модель Word2Vec
class Word2VecModel(Basic):
    direction = None
    code = None
    model = None

    def __init__(self, direction, code):
        super().__init__()
        self.direction = direction
        self.code = code
        self.db = self.connect_db(db_name)

    # Загрузка модели
    def load_model(self):
        try:
            file_path = f"{models_dir}/{self.direction}/{self.code}.model"

            if os.path.exists(file_path):
                self.model = Word2Vec.load(file_path)
                self.logger.debug(f"Model successfully loaded ({self.direction}|{self.code})")
            else:
                self.logger.warning(f"Model file not found ({file_path})")

        except Exception as ex:
            self.logger.error(ex)

    # Создание и тренировка модели
    def create_and_train_model(self, sentences):
        path = f"{models_dir}/{self.direction}"

        try:
            if not os.path.exists(path):
                os.mkdir(path)

            model = Word2Vec(min_count=1, vector_size=2000, sentences=sentences, window=2, negative=10, sample=1e-5,
                             sg=1, workers=4, alpha=0.025, min_alpha=0.0001)
            model.save(path + f"/{self.code}.model")
            self.logger.info(f"Model successfully create, train and save ({path}/{self.code}.model)")

        except Exception as ex:
            self.logger.error(ex)

    # тренировка пользовательской модели
    def create_and_training_user_model(self, foot_prints):
        path = f"{models_dir}/{self.direction}/{self.code}.model"

        sentences = []

        try:
            for fp in foot_prints:
                sentences += fp.sentences

            if not os.path.exists(path):
                model = Word2Vec(min_count=1, vector_size=2000, sentences=sentences, window=2, negative=10, sample=1e-5,
                                 sg=1, workers=4, alpha=0.025, min_alpha=0.0001)
                model.save(path)
                self.model = model
                self.logger.info(f"{self.code} user model was successfully created")
            else:
                self.load_model()
                self.model.train(sentences, total_examples=len(sentences), epochs=10, start_alpha=0.0001,
                                 end_alpha=0.025)
                self.model.save(path)
                self.logger.info(f"{self.code} user model was successfully train")

        except Exception as ex:
            self.logger.error(ex)
            raise Exception(ex)

        # Перемещение отработанных файлов-следов в папку spent
        for fp in foot_prints:
            file_path = f"{foot_prints_dir}/{fp.direction}/{self.code}/{fp.file}"
            if os.path.exists(file_path):
                spent_dir = f"{foot_prints_dir}/{fp.direction}/{self.code}/spent"
                if not os.path.exists(spent_dir):
                    os.mkdir(spent_dir)
                shutil.move(file_path, spent_dir)

    # Сравнение двух моделей
    @staticmethod
    def comparison_models(model, model2, words):
        similarities = 0
        hits = {}

        for word in words.keys():
            key_weight = words[word]['weight']

            m1_top = model.model.wv.most_similar(positive=[word], topn=2000)

            try:
                m2_top = dict(model2.model.wv.most_similar(positive=[word], topn=2000))
            except KeyError:
                continue

            coin_vals_1 = []
            coin_vals_2 = []
            for m in m1_top:
                if m[0] in m2_top:
                    coin_vals_1.append(m[1])
                    coin_vals_2.append(m2_top[m[0]])

            if len(coin_vals_1) == 0:
                continue

            coin_vals_1 = Word2VecModel.centred_and_normalization_vec(coin_vals_1)
            coin_vals_2 = Word2VecModel.centred_and_normalization_vec(coin_vals_2)

            sc = Word2VecModel.similarity_cosine(np.array(coin_vals_1),
                                                 np.array(coin_vals_2))

            if sc < 0 or sc > 1:
                continue

            similarities += sc * key_weight
            hits[word] = {'weight': key_weight, 'sim_cos': sc, 'p': sc * key_weight}

        hits = dict(sorted(hits.items(), key=lambda x: x[1]['p'], reverse=True))
        return similarities, hits

    def get_logger(self):
        super().get_logger()

    def connect_db(self, name):
        return super().connect_db(name)

    @staticmethod
    def similarity_cosine(vec1, vec2):
        cosine_similarity = np.dot(matutils.unitvec(vec1), matutils.unitvec(vec2))
        return cosine_similarity

    @staticmethod
    def centred_and_normalization_vec(vec):
        mean_vec = np.mean(vec)
        vec = list((i - mean_vec for i in vec))
        max_vec = max(vec, key=abs)
        if max_vec == 0:
            return vec
        vec = list((i / max_vec for i in vec))
        return vec


def main():
    """
    # Создание объекта стандарта по направлению
    standard = Standard('it')

    # Сформировать модели по коду ФГОС или ПС
    standard.form_fgos_models('44')
    standard.form_ps_models('06')

    # Сформировать все модели по направлению
    standard.form_direction_models()

    # Записать в JSON файл сравнение всех стандартов для направления
    Standard('ped').comparison_all_models("all")
    Standard('it').comparison_all_models("all")

    # Записать в JSON файл сравнение стандартов с моделью пользователя для направления
    Standard('ped').comparison_all_models("user", "test_user")
    Standard('it').comparison_all_models("user", "test_user")
    # В БД
    Standard('it').comparison_all_models("user", "test_user", None, "db")
    """

    """
    # Получение всех цифровых следов из директории, с извлечением и обработкой предложений
    footprints = FootPrint('test_user').get_all_foot_prints_from_dir()
    # Обучение/дообучение модели пользователя по цифровым следам
    Word2VecModel('users', 'test_user').create_and_training_user_model(footprints)
    
    # Обучение и оценка всех пользователей направления с сохранением результатов в базе данных
    Standard.training_and_comparison_all_users_by_dir("it", 10)
    Standard.training_and_comparison_all_users_by_dir("ped")
    """

    #Standard('it').comparison_all_models("all")
    #Standard('it').comparison_all_models("user", "test_user", None, "db")

    #footprints = FootPrint('test_user', 'it').get_all_foot_prints_from_dir(10)
    #Word2VecModel('users', 'test_user').create_and_training_user_model(footprints)

    Standard.training_and_comparison_all_users_by_dir("it", 10)


if __name__ == '__main__':
    main()
